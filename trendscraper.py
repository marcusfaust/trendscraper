"""

from flask import Flask

app = Flask(__name__)


@app.route('/')
def hello_world():
    return 'Hello World!'


if __name__ == '__main__':
    app.run()
"""

import argparse, requests, re, zipfile, StringIO
import httplib2, urllib
from pptx import Presentation
from oauth2client.client import OAuth2WebServerFlow, flow_from_clientsecrets, AccessTokenCredentials
from oauth2client import tools
from oauth2client.file import Storage
from apiclient.discovery import build
from oauth2client.tools import run_flow
from models import *
from sqlalchemy.orm import sessionmaker
import os
import config
import base64


from apiclient import errors

ftplinks = []

class GmailSession:


    def __init__(self , session):
        self.refresh_token = session.query(RefreshToken).filter_by(id=1).first()


    gmail_client_id = os.environ.get('GMAIL_CLIENT_ID')
    gmail_client_secret = os.environ.get('GMAIL_CLIENT_SECRET')

    """
    #Present in case re-authorization is needed or if Scope changes

    flow = OAuth2WebServerFlow(client_id=gmail_client_id,
                               client_secret=gmail_client_secret,
                               scope='https://www.googleapis.com/auth/gmail.modify',
                               redirect_uri='https://www.example.com/oauth2callback',
                               access_type='offline',
                               approval_prompt='force')

    parser = argparse.ArgumentParser(parents=[tools.argparser])
    flags = parser.parse_args()

    storage = Storage('Trendscraper.json')

    credentials = run_flow(flow, storage, flags)
    print "hello"
    """

    def getAccessToken(self, session):

        tokenURL = "https://www.googleapis.com/oauth2/v3/token"

        params = {"client_id": self.gmail_client_id,
                  "client_secret": self.gmail_client_secret,
                  "refresh_token": str(self.refresh_token.token),
                  "grant_type": "refresh_token"}

        r = requests.post(tokenURL, params=params)
        results = r.json()
        atoken = results['access_token']

        # Store refresh_token
        #session.query(RefreshToken.filter_by(id=1).update(dict(token=rtoken))
        #db.session.commit()

        return atoken

#Function that access a Presentation object, finds Array summary slide within, and scrapes information.  Attemps to call function to populate info into db.
def scrape_pptx(prs, session):
    array_summary_found = False
    system_summary_found = False
    ARRAYSUMMARY = {}
    for slide in prs.slides:
        for shape in slide.shapes:

            if (shape.name == 'Subtitle 2'):
                m = re.search("(\w+.*)\nAs of.*", shape.text)
                if m is not None:
                    ARRAYSUMMARY['Customer'] =  m.group(1).rstrip()

            if (shape.has_text_frame and re.search("^Summary", shape.text)):
                array_summary_found = True
                continue

            if (shape.has_text_frame and shape.text == 'System Summary'):
                system_summary_found = True
                continue

            if (shape.has_table and system_summary_found):
                for row in range(0,len(shape.table.rows._tbl.tr_lst)):
                    ARRAYSUMMARY[shape.table.cell(row,0).text_frame.text] = shape.table.cell(row,1).text_frame.text
                    system_summary_found = False

            if (shape.has_table and array_summary_found):
                #for column in range(0,len(shape.table.rows._tbl.tr_lst)):
                    nametext = shape.table.cell(1,0).text_frame.text
                    if (re.search("\(", nametext)):
                        ARRAYSUMMARY['array_name'], ARRAYSUMMARY['serial_no'] = nametext.split( )
                        ARRAYSUMMARY['serial_no'] =  ARRAYSUMMARY['serial_no'].replace("(", "")
                        ARRAYSUMMARY['serial_no'] =  ARRAYSUMMARY['serial_no'].replace(")", "")
                        array_summary_found = False
                    else:
                        ARRAYSUMMARY['array_name'] = nametext
                        ARRAYSUMMARY['serial_no'] = nametext
                        array_summary_found = False

    #Populate DB
    #print ARRAYSUMMARY
    populate_db_from_scrape(ARRAYSUMMARY, session)
    print "Database Populated"

#Function that takes ARRAYSUMMARY dict and populates db accepting session object also.
def populate_db_from_scrape(arraysummary, session):
    summ = Summary(arraysummary['Customer'],
               arraysummary['array_name'],
               arraysummary['serial_no'],
               arraysummary['FLARE'],
               arraysummary['Configured LUN Capacity'],
               arraysummary['Disk Capacity'],
               arraysummary['% Reads'],
               arraysummary['Front End IOPS - avg'],
               arraysummary['Front End IOPS - 95th'],
               arraysummary['Front End IOPS - max'],
               arraysummary['Model'],
               arraysummary['Avg IO Size (KB)'])

    session.add(summ)
    session.commit()



#engine = create_engine(URL(**config.DATABASE))
engine = create_engine(config.SQLALCHEMY_DATABASE_URI)
Session = sessionmaker(bind=engine)
session = Session()

#TEST
#prs = Presentation('vnx.pptx')
#scrape_pptx(prs, session)

gmailsession = GmailSession(session)
access_token = gmailsession.getAccessToken(session)
new_credentials = AccessTokenCredentials(access_token, '', '')


# Path to the client_secret.json file downloaded from the Developer Console
CLIENT_SECRET_FILE = 'Trendscraper.json'

# Check https://developers.google.com/gmail/api/auth/scopes for all available scopes
OAUTH_SCOPE = 'https://www.googleapis.com/auth/gmail.modify'

# Location of the credentials storage file
#STORAGE = Storage('gmail.storage')

#HTTP object
http = httplib2.Http()
http = new_credentials.authorize(http)

#Gmail Service
gmail_service = build('gmail', 'v1', http=http)

# Retrieve a list of all messages
messages = gmail_service.users().messages().list(userId='me').execute()

#Get Contents of Each Message and Populate ftpdict
if messages.has_key('messages'):
    for message in messages['messages']:
        message_text = gmail_service.users().messages().get(id=message['id'], userId='me', format='full').execute()
        for part in message_text['payload']['parts']:
            print "hello"
            if part.has_key('body') and part['body'].has_key('data'):
                msg = base64.urlsafe_b64decode(part['body']['data'].encode('UTF-8'))
                ftpmatches = re.findall('\r\n(.*assessment_download.*.zip)\r\n', msg)
                for link in ftpmatches:
                    ftplinks.append(link)

        #Delete Email Message
        gmail_service.users().messages().trash(id=message['id'], userId='me').execute()

#Download Each FTP Link and find pptx file in zipfile
for link in ftplinks:
    print "Downloading " + link
    filename = re.findall('\/(\w+.zip)', link)
    file = urllib.urlretrieve(link, "./" + filename[0])
    print "Done"

    #Create zipfile object
    zf = zipfile.ZipFile(filename[0], 'r')
    for archive in zf.filelist:
        if re.match('.*VNX Profile.pptx', archive.filename):
            temp_pptx = zf.read(archive.filename)
            stream = StringIO.StringIO(temp_pptx)
            prs = Presentation(stream)
            scrape_pptx(prs, session)

    #Delete downloaded zip file
    os.remove("./" + filename[0])
